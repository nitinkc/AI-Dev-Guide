#!/usr/bin/env python3
"""
Generate PowerPoint presentation for AI-Powered Development Automation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Define color scheme
DARK_BG = RGBColor(13, 17, 23)      # #0d1117
BLUE_ACCENT = RGBColor(88, 166, 255)  # #58a6ff
LIGHT_TEXT = RGBColor(230, 237, 243)  # #e6edf3
GRAY_TEXT = RGBColor(139, 148, 158)    # #8b949e

def set_slide_background(slide):
    """Set dark background for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title_slide(prs, title, subtitle, date_text):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = BLUE_ACCENT
    title_p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.italic = True
    subtitle_p.font.color.rgb = GRAY_TEXT
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = date_text
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = GRAY_TEXT
    date_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_elements):
    """
    Add a content slide with title and mixed content
    content_elements: list of dicts with 'type', and other fields depending on type
    Types: 'bullet', 'callout', 'stat', 'table', 'text'
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = BLUE_ACCENT

    # Add content
    y_position = 1.4
    for element in content_elements:
        if element['type'] == 'bullet':
            y_position = add_bullet_points(slide, element['items'], y_position, element.get('level', 0))
        elif element['type'] == 'text':
            y_position = add_text_box(slide, element['text'], y_position, element.get('size', 18), element.get('bold', False))
        elif element['type'] == 'stat':
            y_position = add_stat_box(slide, element['text'], y_position)
        elif element['type'] == 'callout':
            y_position = add_callout_box(slide, element['text'], y_position)
        elif element['type'] == 'table':
            y_position = add_table_content(slide, element['headers'], element['rows'], y_position)

    return slide

def add_bullet_points(slide, items, y_pos, level=0):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(Inches(0.8 + (0.3 * level)), Inches(y_pos), Inches(8.4 - (0.3 * level)), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_TEXT
        p.level = level
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.3

    return y_pos + 0.25 + (len(items) * 0.35)

def add_text_box(slide, text, y_pos, size=18, bold=False):
    """Add a text box"""
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = LIGHT_TEXT
    p.space_before = Pt(6)
    p.space_after = Pt(6)

    return y_pos + 0.3 + (size / 72) * 0.2

def add_stat_box(slide, text, y_pos):
    """Add a highlighted stat box"""
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(88, 166, 255) if False else RGBColor(30, 50, 90)
    shape.line.color.rgb = BLUE_ACCENT
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = LIGHT_TEXT
    p.alignment = PP_ALIGN.LEFT

    return y_pos + 1.2

def add_callout_box(slide, text, y_pos):
    """Add a callout box"""
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 50, 90)
    shape.line.color.rgb = BLUE_ACCENT
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.alignment = PP_ALIGN.CENTER

    return y_pos + 1.3

def add_table_content(slide, headers, rows, y_pos):
    """Add table content"""
    # Dimensions
    left = Inches(0.5)
    top = Inches(y_pos)
    width = Inches(9)
    height = Inches(0.4 + len(rows) * 0.4)

    # Add table
    rows_count = len(rows) + 1
    cols_count = len(headers)
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # Set column widths
    for col_idx in range(cols_count):
        table_shape.columns[col_idx].width = width // cols_count

    # Add headers
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 50, 90)
        text_frame = cell.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        p.alignment = PP_ALIGN.CENTER

    # Add rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(20, 30, 50)
            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_TEXT
            p.alignment = PP_ALIGN.CENTER

    return y_pos + 0.5 + (rows_count * 0.4)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Title
    add_title_slide(prs,
        "AI-Powered Development Automation",
        "Automating the manual, accelerating the meaningful",
        "April 2026"
    )

    # SLIDE 2: The Problem
    add_content_slide(prs, "The Problem: Engineering Time Wasted on Toil", [
        {'type': 'bullet', 'items': [
            'Developers spend 30–40% of time on repetitive, low-value tasks',
            'Pain points:',
            '  • JIRA ticket triage, grooming, and hand-off are manual and slow',
            '  • Playwright E2E test failures require 30–120 min/failure to investigate',
            '  • QA cycles bottlenecked by developer availability',
            '  • Spring Boot boilerplate consumes senior engineer time',
        ]},
        {'type': 'stat', 'text': '"Each CI failure costs ~1 hour of developer time. 10 failures/week = 1 dev-week/month lost"'},
    ])

    # SLIDE 3: Status Quo
    add_content_slide(prs, "Status Quo: Manual Workflow Bottlenecks", [
        {'type': 'text', 'text': 'JIRA → Code Flow:', 'bold': True, 'size': 20},
        {'type': 'text', 'text': 'JIRA ticket → reads → writes code → writes tests → opens PR (Days elapsed)'},
        {'type': 'text', 'text': 'CI Failure Flow:', 'bold': True, 'size': 20},
        {'type': 'text', 'text': 'CI fails → investigates → reproduces → fixes → re-pushes (Hours elapsed)'},
        {'type': 'callout', 'text': 'Key insight: These workflows are structured, repeatable, and automatable'},
    ])

    # SLIDE 4: Our Solution
    add_content_slide(prs, "Our Solution: AI Agentic Automation Platform", [
        {'type': 'text', 'text': 'Event-driven AI agent system handling routine development tasks end-to-end', 'bold': False, 'size': 18},
        {'type': 'text', 'text': 'Core Technology Stack:', 'bold': True, 'size': 20},
        {'type': 'bullet', 'items': [
            'LLMs (Claude 3.5 Sonnet / GPT-4o) — reasoning and code generation',
            'RAG — vector-indexed codebase for context-aware retrieval',
            'LangGraph — stateful orchestration with human-in-the-loop gates',
            'MCP Servers — standardised tool integrations for JIRA, GitHub, Playwright',
        ], 'level': 0},
        {'type': 'callout', 'text': '"AI proposes. Humans approve. CI validates."'},
    ])

    # SLIDE 5: Solution Architecture
    add_content_slide(prs, "Solution Architecture (High-Level)", [
        {'type': 'text', 'text': 'Triggers: JIRA webhook, CI failure webhook', 'size': 18},
        {'type': 'text', 'text': 'Orchestration: LangGraph state machine with agent nodes'},
        {'type': 'text', 'text': 'Tool Layer: JIRA MCP, GitHub MCP, Playwright MCP, pgvector code index'},
        {'type': 'text', 'text': 'Outputs: GitHub PR, PR comment with RCA, JIRA comment'},
        {'type': 'bullet', 'items': [
            'Fully event-driven, no polling, resumable via PostgreSQL checkpointer',
            'Key design principle: Human interrupt gate is non-negotiable before PR creation',
        ]},
    ])

    # SLIDE 6: Use Case 1
    add_content_slide(prs, "Use Case 1: JIRA Ticket → Pull Request", [
        {'type': 'text', 'text': 'Trigger: Ticket assigned to AI agent in JIRA', 'bold': True, 'size': 20},
        {'type': 'text', 'text': 'Agent Flow: Read → Classify → RAG retrieve → ReAct plan → Generate → Test → Validate → Human review → Open PR', 'size': 18},
        {'type': 'text', 'text': 'Guardrails: Diff scope validator (>50 lines → human flag), multi-service → escalation', 'size': 18},
        {'type': 'table', 'headers': ['Metric', 'Value'], 'rows': [
            ['End-to-end run time', '3–15 minutes per ticket'],
            ['API cost', '$0.05–$2.00 per ticket'],
            ['Manual time saved', '2–4 hours per ticket'],
        ]},
    ])

    # SLIDE 7: Use Case 2
    add_content_slide(prs, "Use Case 2: Playwright CI Failure → RCA & Auto-Fix", [
        {'type': 'text', 'text': 'Trigger: GitHub Actions detects Playwright test failure', 'bold': True, 'size': 20},
        {'type': 'text', 'text': 'Agent Flow: Download report → Classify failure type → Auto-fix OR RCA document → Push/Post', 'size': 18},
        {'type': 'text', 'text': 'Failure types: Flaky timing, broken selectors, API contract breaks, environment issues', 'size': 18},
        {'type': 'table', 'headers': ['Metric', 'Target'], 'rows': [
            ['Agent analysis time', '< 2 minutes (vs 30–120 manual)'],
            ['Correct classification rate', '>90%'],
            ['Auto-fixes passing on first try', '>70%'],
            ['Investigation time reduction', '>60%'],
        ]},
    ])

    # SLIDE 8: The Advantage
    add_content_slide(prs, "The Advantage: Speed, Accuracy & Developer Focus", [
        {'type': 'table', 'headers': ['Dimension', 'Before', 'After'], 'rows': [
            ['JIRA → PR cycle time', '1–3 days', '3–15 minutes'],
            ['CI failure investigation', '30–120 min', '< 2 minutes'],
            ['Developer context switching', 'High', 'Low'],
            ['Test coverage per PR', 'Variable', 'Enforced'],
            ['Cost per ticket', '—', '$0.05–$2.00'],
        ]},
        {'type': 'text', 'text': 'Human engineers refocus on: architecture decisions, complex features, code review quality'},
    ])

    # SLIDE 9: Safety & Governance
    add_content_slide(prs, "Safety & Governance", [
        {'type': 'bullet', 'items': [
            'Every AI action is reversible or gated before taking effect',
            'Human-in-the-loop via LangGraph interrupt before PR creation and commit push',
            'Input/output guardrails (Guardrails AI) — no PII or secrets in LLM context',
            'Secret scanning (gitleaks/TruffleHog) on generated code before push',
            'Full audit trail: LangSmith/Langfuse tracing, PostgreSQL state checkpointing',
            'AI-generated PRs are clearly labelled — no stealth automation',
        ]},
        {'type': 'callout', 'text': '"The agent cannot merge. Only humans merge."'},
    ])

    # SLIDE 10: Team & Phased Delivery
    add_content_slide(prs, "How We Approached It: Team & Phased Delivery", [
        {'type': 'text', 'text': 'Team Structure (4 roles):', 'bold': True, 'size': 20},
        {'type': 'bullet', 'items': [
            'AI/ML Engineer — LangGraph orchestration, RAG pipeline, LLM integration',
            'Backend Engineer — Spring Boot services, MCP adapters, build sandbox',
            'QA/Automation Engineer — Playwright MCP, test classification logic',
            'DevOps/Platform Engineer — CI webhooks, Docker sandbox, PostgreSQL checkpointer',
        ]},
        {'type': 'text', 'text': 'Phased Timeline:', 'bold': True, 'size': 20},
        {'type': 'table', 'headers': ['Phase', 'Weeks', 'Focus'], 'rows': [
            ['Phase 1', '1–4', 'Foundations: RAG, JIRA MCP, LangGraph, gates'],
            ['Phase 2', '5–8', 'Use Case 1: JIRA → PR (bugs, features)'],
            ['Phase 3', '9–12', 'Use Case 2: Playwright RCA + auto-fix + CI'],
            ['Phase 4', '13+', 'Expansion: code review, docs, dependencies'],
        ]},
    ])

    # SLIDE 11: Expansion Roadmap
    add_content_slide(prs, "What Else This Platform Enables", [
        {'type': 'text', 'text': 'Expansion Roadmap — Same Stack, New Targets:', 'bold': True, 'size': 20},
        {'type': 'bullet', 'items': [
            'Automated code review comments (security, style, coverage) on every PR',
            'Documentation generation from merged code changes',
            'Dependency upgrade PRs triggered by security advisories',
            'On-call runbook execution from PagerDuty alerts',
            'Sprint retrospective summaries from completed JIRA tickets',
            'Test gap detection — weekly scan for uncovered code paths',
        ]},
        {'type': 'callout', 'text': 'The platform is an investment in infrastructure, not a one-off script'},
    ])

    # SLIDE 12: Summary & Next Steps
    add_content_slide(prs, "Summary & Next Steps", [
        {'type': 'text', 'text': 'Three-Sentence Executive Summary:', 'bold': True, 'size': 20},
        {'type': 'bullet', 'items': [
            '1. We built an AI agent platform that turns JIRA tickets into PRs and Playwright failures into root cause reports — with human approval gates',
            '2. The system uses production-grade LLMs, RAG-indexed codebase, and standardised MCP integrations',
            '3. Initial use cases show 10–50× time savings with a clear expansion roadmap',
        ]},
        {'type': 'text', 'text': 'Call to Action:', 'bold': True, 'size': 20},
        {'type': 'bullet', 'items': [
            'Approve Phase 2 resourcing to productionise Use Case 1',
            'Schedule a live demo of the JIRA → PR flow',
            'Define the list of Spring Boot services to onboard first',
        ]},
    ])

    # Save presentation
    output_path = '/Users/sgovinda/Learn/aiDevGuide/AI-Powered-Development-Automation.pptx'
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")

if __name__ == '__main__':
    create_presentation()

